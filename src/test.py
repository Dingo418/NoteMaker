from pptx import Presentation
from pathlib import Path


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    extracted_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        
        # Extract text from slide content
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        # Extract notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
        
        extracted_text.append({
            "slide_number": i + 1,
            "slide_text": "\n".join(slide_text),
            "notes_text": notes_text
        })
    
    return extracted_text

# Example usage
if __name__ == "__main__":
    pptx_file = Path('FCDS.pptx')  # Change to your file path
    text_data = extract_text_from_pptx(pptx_file)
    
    for slide in text_data:
        print(f"Slide {slide['slide_number']}:")
        print("Slide Text:")
        print(slide['slide_text'])
        print("Notes:")
        print(slide['notes_text'])
        print("-" * 40)
