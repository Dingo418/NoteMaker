import os
import yt_dlp
import textract
import requests
from pathlib import Path
from pptx import Presentation

TEMP_HTML_FILE = Path("data/temp.html")
TEMP_WAV_FILE = Path("data/temp_audio")

def read_file_content(file_path : Path) -> str:
    """
    Open and reads a file
    """
    try:
        with open(file_path, 'r') as f:
            responses = f.read()
        return responses
    except Exception as e:
        print(f"An error occured while reading the file: {e}")
        return None

def extract_text_from_pptx(file_path : Path) -> str:
    """
    Extract text on the slide and the notes of a powerpoint
    """
    prs = Presentation(file_path)
    extracted_text = []
  
    # Goes into each slide and rips all the text out
    for i, slide in enumerate(prs.slides):
        slide_text = []
      
        # Extract text from slide content
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
      
        # Extract notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
      
        extracted_text.append({
            "slide_number": i + 1,
            "slide_text": "\n".join(slide_text),
            "notes_text": notes_text
        })
  
    return str(extracted_text) # Returns it as string

def get_youtube(URL : str) -> None:
    """
    Get's the youtube video, and returns it as wav
    """
    ydl_opts = {
        'format': 'm4a/bestaudio/best',
        'postprocessors': [{  
            'key': 'FFmpegExtractAudio',
            'preferredcodec': 'wav',
        }],
        'outtmpl': str(TEMP_WAV_FILE),
        'quiet': True,       # Suppresses output logs
        'no_warnings': True  # Suppresses warnings
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        error_code = ydl.download(URL)
    # If the rip is not successful, raise a exception with the error code
    if error_code != 0:
        raise ValueError("youtube-dl error: ", error_code)
    
def fetch_website(url : str) -> None:
    """
    Grabs a website's HTML and saves it to a file.
    """
    try:
        response = requests.get(url)
        response.raise_for_status()

    except requests.exceptions.RequestException as e:
        print(f"Error fetching URL: {e}")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
    
    with open(TEMP_HTML_FILE, 'a') as file:
        file.write(response.text)

def delete_temp_files() -> None:
    """
    Remove any files if they exist
    """
    #Checks if path is a file, then delete
    if os.path.isfile(TEMP_HTML_FILE):
        os.remove(TEMP_HTML_FILE)
    if os.path.isfile(TEMP_WAV_FILE):
        os.remove(TEMP_WAV_FILE)

def extract_text_from_file(file_path : str) -> str:
    """
    Extracts text from various file types.
    """
    # Checks if file_path is a website
    if "http" in file_path:
        # youtube is special so it gets it's own module
        if "youtube.com" in file_path:
            youtube_url = file_path
            get_youtube(youtube_url)
            print("Youtube video has been converted to wav:")
            file_path = Path("data/temp_audio.wav")
        else:
           # Otherwise just rip it
           fetch_website(file_path)
           file_path = TEMP_HTML_FILE
    
    file_path = Path(file_path)
    file_extension = file_path.suffix
    if file_extension in (".csv", ".doc", ".docx", ".eml", ".epub", ".gif", ".htm", ".html", ".jpeg", ".jpg", ".json", ".log", ".mp3", ".msg", ".odt", ".ogg", ".pdf", ".png", ".ps", ".psv", ".rtf", ".tab", " tff", ".tif", ".tiff", ".tsv", " .txt", ".wav", ".xls", ".xlsx"):
        text = textract.process(file_path).decode()
    elif file_extension in (".txt", ".md"):
        text =  read_file_content(file_path)
    elif file_extension == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {file_extension}")
    
    delete_temp_files()
    return text